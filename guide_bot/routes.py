# guide_bot/routes.py
import getpass
import io
import re
from flask import Blueprint, abort, render_template, redirect, send_file, url_for, flash, request, session
from langchain_chroma import Chroma
import markdown
from guide_bot.models import Document
from guide_bot.forms import DocumentFileForm, DocumentFolderForm
from app import db
from PyPDF2 import PdfReader
import docx2txt
from pptx import Presentation
import pandas
from langchain_huggingface import HuggingFaceEmbeddings
from dotenv import load_dotenv
import os
import bcrypt
from langchain_core.documents import Document as ChatDocument
from werkzeug.utils import secure_filename
from guide_bot import conversation_chat, create_conversational_chain, load_vector_store, save_uploaded_file, load_saved_files, split_documents

guide_bot = Blueprint('guide_bot', __name__)

GOOGLE_API_KEY=os.getenv('GOOGLE_API_KEY')

# Don't forget to update the manage_documents route to save documents with their ids
@guide_bot.route('/guide-bot/documents', methods=['GET', 'POST'])
def manage_documents():
    file_form = DocumentFileForm()
    folder_form = DocumentFolderForm()
    
    # Ambil page dan items_per_page dari parameter query, default page=1 dan per_page=10
    page = request.args.get('page', 1, type=int)
    items_per_page = request.args.get('items_per_page', 10, type=int)
    
    # Ambil query pencarian dari parameter query
    search_query = request.args.get('search_query', '')
    
    # Filter dokumen berdasarkan search query jika ada
    if search_query:
        documents = Document.query.filter(Document.title.ilike(f'%{search_query}%')).paginate(page=page, per_page=items_per_page)
    else:
        documents = Document.query.paginate(page=page, per_page=items_per_page)
    
    # Proses penambahan dokumen
    if file_form.validate_on_submit() or folder_form.validate_on_submit():
        files = request.files.getlist('files')
        for file in files:
            filename = secure_filename(file.filename)
            # Simpan dokumen dalam database
            new_document = Document(title=filename, file=file.read())
            db.session.add(new_document)
            db.session.commit()

            # Save the file to disk with its document ID
            save_dir = 'data'
            if not os.path.exists(save_dir):
                os.makedirs(save_dir)
            file_path = os.path.join(save_dir, f"{filename}")
            with open(file_path, 'wb') as f:
                f.write(new_document.file)

        flash('Document added successfully!')
        return redirect(url_for('guide_bot.manage_documents', page=page, items_per_page=items_per_page, search_query=search_query))
    
    # Render halaman dengan form dan dokumen terpaginate
    return render_template('guide_bot/manage_documents.html', 
                           file_form=file_form, 
                           folder_form=folder_form, 
                           documents=documents, 
                           items_per_page=items_per_page, 
                           search_query=search_query)
    
@guide_bot.route('/guide-bot/documents/delete-multiple', methods=['POST'])
def delete_multiple_documents():
    document_ids = request.form.getlist('document_ids')  # Use getlist to handle multiple values

    if not document_ids:
        flash('No documents selected for deletion.', 'warning')
        return redirect(url_for('guide_bot.manage_documents'))

    deleted_files = []
    not_found_files = []
    errors = []

    for doc_id in document_ids:
        if doc_id:  # Ensure id is not empty
            document = Document.query.get(doc_id)
            if document:
                db.session.delete(document)
                file_path = os.path.join('data', f"{document.id}_{secure_filename(document.title)}")
                if os.path.exists(file_path):
                    try:
                        os.remove(file_path)
                        deleted_files.append(file_path)
                    except Exception as e:
                        errors.append(f'Error removing file {file_path}: {e}')
                else:
                    not_found_files.append(file_path)
            else:
                flash(f'Document with ID {doc_id} not found.', 'warning')
    
    db.session.commit()

    if deleted_files:
        flash(f'Selected documents deleted successfully: {", ".join(deleted_files)}', 'success')
    if not_found_files:
        flash(f'Files not found: {", ".join(not_found_files)}', 'warning')
    if errors:
        flash(f'Errors occurred: {", ".join(errors)}', 'error')

    return redirect(url_for('guide_bot.manage_documents'))

@guide_bot.route('/guide-bot/documents/delete/<int:id>', methods=['POST'])
def delete_document(id):
    document = Document.query.get_or_404(id)
    
    # Delete the document from the database
    db.session.delete(document)
    db.session.commit()

    # delete the file from the server
    file_path = os.path.join('data', secure_filename(document.title))
    if os.path.exists(file_path):
        try:
            os.remove(file_path)
            print(f"Successfully deleted file: {file_path}")  # Debugging statement
        except Exception as e:
            print(f"Error removing file {file_path}: {e}")
    else:
        print(f"File not found: {file_path}")

    flash('Document deleted successfully!')
    return redirect(url_for('guide_bot.manage_documents'))

@guide_bot.route('/guide-bot/documents/view/<int:id>')
def view_document(id):
    document = Document.query.get_or_404(id)
    return render_template('guide_bot/view_document.html', document=document)

@guide_bot.route('/guide-bot/document/file/<int:document_id>')
def get_document_file(document_id):
    document = Document.query.get_or_404(document_id)
    if document.file:
        # Set appropriate MIME type based on file type
        return send_file(io.BytesIO(document.file), mimetype='application/pdf', as_attachment=False)
    else:
        abort(404)

@guide_bot.route('/guide-bot/documents/download/<int:id>')
def download_document(id):
    document = Document.query.get_or_404(id)
    return send_file(io.BytesIO(document.file), as_attachment=True, download_name=document.title)

# Initialize variables
history = []
generated = ["Selamat datang di GuideBot! Tanyakan sesuatu pada saya 😊️"]
past = []

@guide_bot.route('/guide-bot/chat', methods=['GET', 'POST'])
def chat():
    global history, generated, past

    if request.method == 'POST':
        user_input = request.form['user_input']        
        vector_store = load_vector_store(embeddings)
        print("Vector store loaded")
        if vector_store:
            print("Vector store is not empty")
            chain = create_conversational_chain(vector_store)
            print("Conversational chain created")
            print("History:", history)
            output, source_documents = conversation_chat(user_input, chain, history)
            print("Output:", output)
            print("Conversation completed"  )
            past.append(user_input)
            generated.append({
                "message": markdown.markdown(output),
                "source_documents": source_documents
            })
        else:
            flash('⚠️ No documents uploaded by admin yet.', 'warning')

    return render_template('guide_bot/index.html', past=past, generated=generated)

embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2", model_kwargs={'device': 'cpu'})

@guide_bot.route('/guide-bot/reload-vector-db', methods=['GET'])
def reload_vector_db():
    vector_store = Chroma(
        collection_name="SPIL",
        embedding_function=embeddings or HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2", model_kwargs={'device': 'cpu'}), 
        persist_directory="vector_store"
    )
    # Load all documents from the SQL database
    documents = Document.query.all()
    document_ids = [str(document.id) for document in documents]    
    vector_store_ids = vector_store.get()['ids']    
    print("Documents:", document_ids)    
    print("Vector store IDs:", vector_store_ids)

    # Prepare documents for vector store
    for document in documents:
        vector_store_ids = vector_store.get()['ids']
        print("Vector store IDs:", vector_store_ids)
        # Check if document already exists in the vector store
        if str(document.id) in vector_store_ids:
            continue           

        file_path = os.path.join('data', secure_filename(document.title))
        
        # Extract text based on file type
        all_text = ""
        if file_path.endswith('.pdf'):
            pdf_reader = PdfReader(file_path)
            for page in pdf_reader.pages:
                text = page.extract_text()
                if text:
                    all_text += text
        elif file_path.endswith('.docx'):
            all_text = docx2txt.process(file_path)        
        elif file_path.endswith('.ods') or file_path.endswith('.xls') or file_path.endswith('.xlsx'):
            df = pandas.read_excel(file_path)
            all_text = df.to_string()
        elif file_path.endswith('.pptx'):
            presentation = Presentation(file_path)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        all_text += shape.text        
        elif file_path.endswith('.txt'):
            with open(file_path, 'r', encoding='utf-8') as file:
                all_text = file.read()
        
        if not all_text.strip():
            print("Failed to extract text from the document.")
            flash("Failed to extract text from the document.")
            continue

        document_obj = ChatDocument(
             page_content= all_text,
             metadata = {
                    "title": document.title,
                    "file_path": file_path,
                    "id": document.id
                }
        )
        
        if len(vector_store_ids) == 0:
            print("Vector store is empty")
            vector_store = Chroma.from_documents(
                embedding=embeddings,
                collection_name="SPIL",
                documents=[document_obj], 
                persist_directory="vector_store",
                ids=[str(document.id)])
            print("Vector store created")
        else:
            print("Vector store is not empty")
            vector_store.add_documents(
                documents=[document_obj], 
                ids=[str(document.id)]
            )
            print("Document with ID", document.id, "added to vector store")

    for vector_id in vector_store_ids:
        if vector_id not in document_ids:
            vector_store.delete(ids=[vector_id])
            print(f"Deleted vector with id {vector_id}")

    flash('Vector database reloaded successfully!')
    return redirect(url_for('guide_bot.manage_documents'))






